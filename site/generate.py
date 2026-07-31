"""
Build script for the ANTH 102 study site.

Reads the Canvas content export (course-data.js + the exported files),
extracts text from the lecture decks, and renders a static site into
../docs using Jinja2 templates.

Run: python generate.py
"""
import json
import os
import re
import shutil
from pathlib import Path

from jinja2 import Environment, FileSystemLoader

SITE_DIR = Path(__file__).resolve().parent
BASE_DIR = SITE_DIR.parent
COURSE_ROOT = BASE_DIR / "Course content"
EXPORT_DIR = COURSE_ROOT / "Summer-2026---ANTH-102-30280--2026-Jul-31_06-31-22-283"
FILES_DIR = EXPORT_DIR / "viewer" / "files"
COURSE_DATA_JS = EXPORT_DIR / "viewer" / "course-data.js"
EXTERNAL_LINKS_TXT = COURSE_ROOT / "ExternalLinks.txt"
FLASHCARDS_PDF = COURSE_ROOT / "Flashcards.pdf"

TEMPLATES_DIR = SITE_DIR / "templates"
STATIC_DIR = SITE_DIR / "static"
OUT_DIR = BASE_DIR / "docs"

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif"}

# Administrative / non-lecture modules dropped from the site entirely per
# user request ("too crowded" -- this is a lecture study guide, not a course
# housekeeping mirror). "Human Anatomy Links" is dropped too: both its
# attachments reference images that no longer exist anywhere in the export,
# so it has zero renderable content.
CUT_MODULES = {
    "Textbook Information",
    "Meet Your Professor",
    "Study Guide",
    "Class Project Assignment",
    "Human Anatomy Links",
    "The Professor's Stupid Music References",
}

# Contiguous position ranges (0-indexed, end-exclusive) over the *kept*
# modules in course order, used purely to group the sidebar/homepage into
# labeled sections. Course order itself is never reshuffled.
ERAS = [
    (0, 2, "Foundations"),
    (2, 6, "Genetics & Evolution"),
    (6, 8, "Human Variation & Forensics"),
    (8, 12, "Primates & Paleoanthropology"),
    (12, 17, "Human Origins"),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def slugify(text):
    text = text.lower()
    text = re.sub(r"\.pptx$", "", text)
    text = re.sub(r"[^a-z0-9]+", "-", text)
    return text.strip("-") or "item"


def unique_slug(base_slug, used):
    slug = base_slug
    n = 2
    while slug in used:
        slug = f"{base_slug}-{n}"
        n += 1
    used.add(slug)
    return slug


def load_course_data():
    raw = COURSE_DATA_JS.read_text(encoding="utf-8").strip()
    raw = re.sub(r"^window\.COURSE_DATA\s*=\s*", "", raw)
    raw = re.sub(r";\s*$", "", raw)
    return json.loads(raw)


def parse_external_links():
    links = []
    if not EXTERNAL_LINKS_TXT.exists():
        return links
    for line in EXTERNAL_LINKS_TXT.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        m = re.match(r"^(\S+)\s*(.*)$", line)
        if not m:
            continue
        url, label = m.group(1), m.group(2).strip()
        links.append({"url": url, "label": label or url})
    return links


def resolve_attachment(rel_path):
    """Return an existing absolute path for an export-relative attachment
    path, falling back to a flat lookup in FILES_DIR by basename (legacy
    nested export paths often point at a file that's actually flat in
    FILES_DIR under the same name). Returns None if nothing matches."""
    abs_path = EXPORT_DIR / rel_path
    if abs_path.exists():
        return abs_path
    fallback = FILES_DIR / os.path.basename(rel_path)
    if fallback.exists():
        return fallback
    return None


_TITLE_CRUFT_PATTERNS = [
    r"\s*\[Autosaved\]\s*$",
    r"\s*\.ppt\s*\(\d+\)\s*$",
    r"\s*\(\d+\)\s*$",
    r"\s*[-–]?\s*updated[\s\-]*\d*\s*$",
    r"\s*[-–]?\s*(Fall|Spring|Summer|Winter)\s*\d{4}\s*$",
    r"\s*[-–]?\s*\d{4}\s*$",
]


def clean_lecture_title(raw_title):
    """Strip the Canvas-era filename cruft (autosave tags, "Updated",
    trailing semester/year, stray version numbers) down to a clean title,
    since the L-code already conveys ordering/version."""
    title = re.sub(r"\.pptx$", "", raw_title, flags=re.IGNORECASE)
    title = re.sub(r"^Lecture\s*\d+\s*[-:]?\s*", "", title, flags=re.IGNORECASE)
    title = re.sub(r"\blecture\b", "", title, flags=re.IGNORECASE)
    changed = True
    while changed:
        changed = False
        for pattern in _TITLE_CRUFT_PATTERNS:
            new_title = re.sub(pattern, "", title, flags=re.IGNORECASE)
            if new_title != title:
                title = new_title
                changed = True
    title = re.sub(r"\s{2,}", " ", title).strip(" -–")
    return title or raw_title


# ---------------------------------------------------------------------------
# Content extraction
# ---------------------------------------------------------------------------

def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(line)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            lines.append(cell.text.strip())
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        # Keep every slide, even image-only ones with no extractable text --
        # the rendered slide image (see export_slide_images) still carries
        # the content, so it still deserves a section on the page.
        slides.append({"number": i, "lines": lines, "notes": notes})
    return slides


def export_slide_images(ppt_app, pptx_path, out_dir):
    """Render every slide of a deck to a real JPG via PowerPoint COM
    automation (ppt_app is a single PowerPoint.Application reused across the
    whole build -- opening/closing per deck is what's slow, not the export
    itself). Returns True on success, False if this deck couldn't be
    rendered (caller falls back to text-only for that lecture)."""
    out_dir.mkdir(parents=True, exist_ok=True)
    pres = None
    try:
        pres = ppt_app.Presentations.Open(str(pptx_path), True, True, False)
        pres.SaveAs(str(out_dir), 17)  # 17 = ppSaveAsJPG
        return True
    except Exception as e:
        print(f"WARNING: slide image export failed for {pptx_path.name}: {e}")
        return False
    finally:
        if pres is not None:
            pres.Close()


_FLASHCARD_NUM_RE = re.compile(r"^\d+\.$")
_FLASHCARD_GLUED_RE = re.compile(r"^([A-Za-z]+-)(\d\)[A-Za-z].*)$")


def _join_flashcard_words(words):
    text = ""
    for w in words:
        if text.endswith("-"):
            text = text[:-1] + w
        elif text:
            text += " " + w
        else:
            text = w
    return text


def extract_flashcards(path):
    """Parse a combined Quizlet PDF export. These render as two columns
    (question left, answer right) per numbered card; plain text extraction
    interleaves the columns, so this reads word bounding boxes instead and
    buckets by x-position. Column split (180) and header/footer bands
    (top<60, top>750) were measured directly off this export's layout."""
    import pdfplumber

    cards = []
    cur_num, cur_q, cur_a = None, [], []

    def flush():
        nonlocal cur_num, cur_q, cur_a
        if cur_num is not None:
            q = _join_flashcard_words(cur_q).strip()
            a = _join_flashcard_words(cur_a).strip()
            if q or a:
                cards.append({"num": cur_num, "q": q, "a": a})
        cur_q, cur_a = [], []

    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            for w in page.extract_words():
                top, x0, text = w["top"], w["x0"], w["text"]
                if top < 60 or top > 750:
                    continue
                if x0 < 30 and _FLASHCARD_NUM_RE.match(text):
                    flush()
                    cur_num = int(text.rstrip("."))
                    continue
                m = _FLASHCARD_GLUED_RE.match(text)
                if x0 < 180 and m:
                    cur_q.append(m.group(1))
                    cur_a.append(m.group(2))
                else:
                    (cur_q if x0 < 180 else cur_a).append(text)
    flush()
    return cards


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def era_for_position(pos):
    for start, end, label in ERAS:
        if start <= pos < end:
            return label
    return "Other"


def build():
    data = load_course_data()
    modules_raw = [m for m in data.get("modules", []) if m["name"] not in CUT_MODULES]
    print(f"Loaded course-data.js: {len(data.get('modules', []))} modules "
          f"({len(modules_raw)} kept after dropping non-lecture modules)")

    used_slugs = set()
    lectures = []
    seen_files = {}  # resolved abs path -> slug, so repeated refs share a page
    stats = {"lecture": 0, "quiz_only": 0, "image": 0, "slide_images": 0}
    unmatched_pptx = []

    ppt_app = None
    try:
        import win32com.client
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    except Exception as e:
        print(f"NOTE: PowerPoint COM unavailable ({e}) -- slides will be text-only.")

    try:
        _build_modules(modules_raw, used_slugs, lectures, seen_files, stats,
                        unmatched_pptx, ppt_app)
    finally:
        if ppt_app is not None:
            ppt_app.Quit()

    external_links = parse_external_links()

    flashcards = []
    if FLASHCARDS_PDF.exists():
        flashcards = extract_flashcards(FLASHCARDS_PDF)
        blank = [c["num"] for c in flashcards if not c["a"]]
        print(f"Flashcards: {len(flashcards)} cards parsed from {FLASHCARDS_PDF.name}"
              + (f" (blank answer in source PDF: {blank})" if blank else ""))

    print(f"Lectures generated: {stats['lecture']} | Images: {stats['image']} | "
          f"Slide images rendered: {stats['slide_images']} decks | "
          f"Kept modules with no pptx (quiz-only, dropped from nav): {stats['quiz_only']}")
    if unmatched_pptx:
        print(f"WARNING: pptx referenced but file not found, skipped: {unmatched_pptx}")

    render(lectures, external_links, flashcards)


def _build_modules(modules_raw, used_slugs, lectures, seen_files, stats, unmatched_pptx, ppt_app):
    for pos, mod in enumerate(modules_raw):
        era = era_for_position(pos)
        module_num = pos + 1

        # Gather this module's supplementary links, quiz, and images once;
        # every lecture page generated from this module shares them.
        links = []
        quiz = None
        images = []
        pptx_items = []

        for item in mod.get("items", []):
            itype = item.get("type")
            if itype == "ExternalUrl":
                links.append({"title": item.get("title", ""), "href": item.get("content", "#")})
            elif itype == "Quizzes::Quiz" and quiz is None:
                quiz = {
                    "title": item.get("title", ""),
                    "points": item.get("pointsPossible"),
                    "question_count": item.get("questionCount"),
                }
            elif itype == "Attachment":
                rel_path = item.get("content", "")
                ext = os.path.splitext(rel_path)[1].lower()
                if ext == ".pptx":
                    pptx_items.append(item)
                elif ext in IMAGE_EXTS:
                    abs_path = resolve_attachment(rel_path)
                    if abs_path:
                        safe_name = re.sub(r"[^A-Za-z0-9._-]", "_", os.path.basename(rel_path))
                        dest = STATIC_DIR / "images" / safe_name
                        dest.parent.mkdir(parents=True, exist_ok=True)
                        shutil.copyfile(abs_path, dest)
                        images.append({"title": item.get("title", ""), "src": f"images/{safe_name}"})
                        stats["image"] += 1

        multi = len(pptx_items) > 1
        for sub_index, item in enumerate(pptx_items, start=1):
            title = item.get("title", "Untitled")
            rel_path = item.get("content", "")
            abs_path = resolve_attachment(rel_path)
            if not abs_path:
                unmatched_pptx.append(title)
                continue

            file_key = str(abs_path.resolve())
            code = f"L{module_num}" + (f".{sub_index}" if multi else "")

            if file_key in seen_files:
                # Same underlying deck referenced again elsewhere (e.g. an
                # addendum listed in two modules) -- point at the existing
                # page rather than re-extracting/duplicating it.
                existing = next(l for l in lectures if l["slug"] == seen_files[file_key])
                lectures.append({**existing, "code": code, "module_name": mod["name"], "era": era})
                stats["lecture"] += 1
                continue

            slug = unique_slug(slugify(title), used_slugs)
            slides = extract_pptx(abs_path)

            slide_image_dir = None
            if ppt_app is not None:
                img_dir = STATIC_DIR / "images" / "lectures" / slug
                if export_slide_images(ppt_app, abs_path, img_dir):
                    slide_image_dir = f"images/lectures/{slug}"
                    stats["slide_images"] += 1

            lecture = {
                "slug": slug,
                "code": code,
                "title": clean_lecture_title(title),
                "module_name": mod["name"],
                "era": era,
                "slides": slides,
                "slide_image_dir": slide_image_dir,
                "links": links,
                "quiz": quiz,
                "images": images,
                "source_file": os.path.basename(rel_path),
            }
            seen_files[file_key] = slug
            lectures.append(lecture)
            stats["lecture"] += 1

        if not pptx_items and quiz:
            stats["quiz_only"] += 1


def render(lectures, external_links, flashcards):
    if OUT_DIR.exists():
        shutil.rmtree(OUT_DIR)
    OUT_DIR.mkdir(parents=True)

    env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)), autoescape=True)

    # Sidebar/homepage nav entries, deduplicated by slug (a shared addendum
    # deck still gets one nav row per code it's filed under).
    nav_lectures = [{
        "code": l["code"], "title": l["title"], "slug": l["slug"], "era": l["era"],
        "cover": f"{l['slide_image_dir']}/Slide1.JPG" if l.get("slide_image_dir") else None,
    } for l in lectures]

    def write(rel_path, template_name, **ctx):
        out_path = OUT_DIR / rel_path
        depth = len(Path(rel_path).parent.parts)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        tmpl = env.get_template(template_name)
        out_path.write_text(tmpl.render(nav_lectures=nav_lectures, depth=depth, **ctx),
                             encoding="utf-8")

    write("index.html", "index.html")

    written_slugs = set()
    for lec in lectures:
        if lec["slug"] in written_slugs:
            continue
        written_slugs.add(lec["slug"])
        write(f"lectures/{lec['slug']}/index.html", "lecture.html", lecture=lec)

    write("resources/index.html", "resources.html", external_links=external_links)

    if flashcards:
        write("flashcards/index.html", "flashcards.html", flashcards=flashcards)

    write_search_index(lectures, flashcards)

    # Unlisted, not private: reachable by anyone with the exact URL, but
    # not crawled/indexed and never linked from anywhere else.
    (OUT_DIR / "robots.txt").write_text("User-agent: *\nDisallow: /\n", encoding="utf-8")

    if STATIC_DIR.exists():
        shutil.copytree(STATIC_DIR, OUT_DIR / "static", dirs_exist_ok=True)

    print(f"Rendered site to {OUT_DIR} ({len(written_slugs)} lecture pages)")


def write_search_index(lectures, flashcards):
    """One chunk per slide with real text, deduplicated by lecture page (a
    deck filed under two codes, like the Homo naledi addendum, is indexed
    once under its first code so search results aren't doubled), plus one
    chunk per flashcard."""
    chunks = []
    indexed_slugs = set()
    for lec in lectures:
        if lec["slug"] in indexed_slugs:
            continue
        indexed_slugs.add(lec["slug"])
        for slide in lec["slides"]:
            text = " ".join(slide["lines"])
            if slide["notes"]:
                text = f"{text} {slide['notes']}"
            text = text.strip()
            if not text:
                continue
            chunks.append({
                "kind": "lecture", "code": lec["code"], "slug": lec["slug"],
                "title": lec["title"], "slide": slide["number"], "text": text,
            })

    for card in flashcards:
        text = f"{card['q']} {card['a']}".strip()
        if not text:
            continue
        chunks.append({
            "kind": "flashcard", "code": "FC", "title": card["q"] or f"Flashcard #{card['num']}",
            "card_num": card["num"], "text": text,
        })

    js_dir = STATIC_DIR / "js"
    js_dir.mkdir(parents=True, exist_ok=True)
    (js_dir / "search-data.js").write_text(
        "const SEARCH_DATA = " + json.dumps(chunks, ensure_ascii=False) + ";\n",
        encoding="utf-8",
    )
    print(f"Search index: {len(chunks)} chunks ({len(indexed_slugs)} lectures, {len(flashcards)} flashcards)")


if __name__ == "__main__":
    build()
